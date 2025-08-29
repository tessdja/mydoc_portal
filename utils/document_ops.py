from __future__ import annotations
from pathlib import Path
from typing import Iterable, List
from fastapi import UploadFile
from langchain.schema import Document
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from logger import GLOBAL_LOGGER as log
from exception.custom_exception import DocumentPortalException

from pptx import Presentation
import pandas as pd
import re

# ingestion gate
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx", ".xlsx", ".csv", ".md"}

def _pptx_to_documents(p: Path) -> List[Document]:
    """
    Extract text (incl. tables) from .pptx using python-pptx.
    """
    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        bits =[]
        for shape in slide.shapes:
            # Text in shapes/placeholders
            if hasattr(shape, "text") and shape.text:
                bits.append(shape.text)
            # Table text
            if getattr(shape, "has_table", False):
                rows = []
                for r in shape.table.rows:
                    rows.append("|".join(c.text for c in r.cells))
                if rows:
                    bits.append("\n".join(rows))

        if bits:
            parts.append(f"---Slide {i} ---\n" + "\n".join(bits))
    content = "\n\n".join(parts) if parts else ""
    return [Document(page_content=content, metadata={"source": str(p), "type": "pptx"})]

def _csv_to_documents(p: Path) -> List[Document]:
    """
    Read CSV and return a single Document with CSV text.
    """
    try:
        df = pd.read_csv(p)
        text = df.to_csv(index=False)
        return [Document(page_content=text, metadata={"source": str(p), "type": "csv", "rows": int(len(df))})]
    except Exception as e:
        raise DocumentPortalException(f"Error reading CSV: {p.name}", e) from e
    
def _xlsx_to_documents(p: Path) -> List[Document]:
    """
    Read each Excel sheet as one Document (sheet-level granularity).
    """
    try:
        xls = pd.ExcelFile(p)
        docs: List[Document] = []
        for sheet in xls.sheet_names:
            df = xls.parse(sheet)
            text = df.to_csv(index=False)
            docs.append(
                Document(
                    page_content=text,
                    metadata = {"source": str(p), "type": "xlsx", "sheet":sheet, "rows": int(len(df))}
                )
            )
        return docs
    except Exception as e:
        raise DocumentPortalException(f"Error reading Excel: {p.name}", e) from e
    
def _md_to_documents(p: Path) -> List[Document]:
    """
    Read Markdown as plain text (preserve headings/code fences).
    """
    try:
        text = p.read_text(encoding="utf-8", errors="ignore")
        # Strip YAML front-matter if present (--- ...--- at top)
        if text.lstrip().startswith("---"):
            text = re.sub(r"^---[\s\S]*?---\s*", "", text, count=1, flags=re.MULTILINE)
        return [Document(page_content=text, metadata={"source": str(p), "type": "md"})]
    except Exception as e:
        raise DocumentPortalException(f"Error reading Markdown: {p.name}", e) from e


def load_documents(paths: Iterable[Path]) -> List[Document]:
    """Load docs using appropriate loader based on extension."""
    docs: List[Document] = []
    try:
        for p in paths:
            ext = p.suffix.lower()
            if ext == ".pdf":
                loader = PyPDFLoader(str(p))
                docs.extend(loader.load())
            elif ext == ".docx":
                loader = Docx2txtLoader(str(p))
                docs.extend(loader.load())
            elif ext == ".txt":
                loader = TextLoader(str(p), encoding="utf-8")
                docs.extend(loader.load())
            elif ext == ".pptx":
                docs.extend(_pptx_to_documents(p))
            elif ext == ".xlsx":
                docs.extend(_xlsx_to_documents(p))
            elif ext == ".md":
                docs.extend(_md_to_documents(p))
            else:
                log.warning("Unsupported extension skipped", path=str(p))
                continue
            
        log.info("Documents loaded", count=len(docs))
        return docs
    except Exception as e:
        log.error("Failed loading documents", error=str(e))
        raise DocumentPortalException("Error loading documents", e) from e

def concat_for_analysis(docs: List[Document]) -> str:
    parts = []
    for d in docs:
        src = d.metadata.get("source") or d.metadata.get("file_path") or "unknown"
        parts.append(f"\n--- SOURCE: {src} ---\n{d.page_content}")
    return "\n".join(parts)

def concat_for_comparison(ref_docs: List[Document], act_docs: List[Document]) -> str:
    left = concat_for_analysis(ref_docs)
    right = concat_for_analysis(act_docs)
    return f"<<REFERENCE_DOCUMENTS>>\n{left}\n\n<<ACTUAL_DOCUMENTS>>\n{right}"

# ---------- Helpers ----------
class FastAPIFileAdapter:
    """Adapt FastAPI UploadFile -> .name + .getbuffer() API"""
    def __init__(self, uf: UploadFile):
        self._uf = uf
        self.name = uf.filename
    def getbuffer(self) -> bytes:
        self._uf.file.seek(0)
        return self._uf.file.read()

def read_pdf_via_handler(handler, path: str) -> str:
    if hasattr(handler, "read_pdf"):
        return handler.read_pdf(path)  # type: ignore
    if hasattr(handler, "read_"):
        return handler.read_(path)  # type: ignore
    raise RuntimeError("DocHandler has neither read_pdf nor read_ method.")